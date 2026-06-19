"""
DocumentProcessor — handles document parsing, recursive chunking, embedding generation,
indexing into Qdrant, and status updates in PostgreSQL.

Implements Requirements: 2.1, 2.2, 2.3, 2.4, 2.7, 2.8, 2.9, 2.10, 2.11, 9.1–9.10.
"""

import asyncio
import io
import logging
import os
import uuid
from typing import List, Optional, Dict, Any

from minio import Minio
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams, PointStruct
from sqlalchemy import text
from llama_index.core.node_parser import SentenceSplitter
from llama_index.core import Document as LlamaDocument

from app.config.configuration_manager import ConfigurationManager
from app.handlers.embedding_generator import EmbeddingGenerator, EmbeddingFailureError
from app.db.connection import get_db

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Processes uploaded documents: parses text, generates chunks/embeddings, and indexes in Qdrant.

    Parameters
    ----------
    config_manager : ConfigurationManager, optional
        Configuration manager instance.
    minio_client : Minio, optional
        MinIO client instance.
    qdrant_client : QdrantClient, optional
        Qdrant client instance.
    embedding_generator : EmbeddingGenerator, optional
        Embedding generator service instance.
    """

    MINIO_HOST = os.getenv("MINIO_HOST", "minio")
    MINIO_PORT = os.getenv("MINIO_PORT", "9000")
    MINIO_ACCESS_KEY = os.getenv("MINIO_ROOT_USER", "minioadmin")
    MINIO_SECRET_KEY = os.getenv("MINIO_ROOT_PASSWORD", "minioadmin")
    MINIO_BUCKET = "documents"

    QDRANT_HOST = os.getenv("QDRANT_HOST", "qdrant")
    QDRANT_PORT = int(os.getenv("QDRANT_PORT", "6333"))
    QDRANT_COLLECTION = "document_chunks"

    def __init__(
        self,
        config_manager: Optional[ConfigurationManager] = None,
        minio_client: Optional[Minio] = None,
        qdrant_client: Optional[QdrantClient] = None,
        embedding_generator: Optional[EmbeddingGenerator] = None,
        db_session=None,
    ) -> None:
        self.config_manager = config_manager
        self._minio_client = minio_client
        self._qdrant_client = qdrant_client
        self._embedding_generator = embedding_generator
        self._db_session = db_session

    async def _execute_db(self, func):
        """Helper to run database operations on injected or direct sessions."""
        if self._db_session is not None:
            return await func(self._db_session)
        else:
            async for session in get_db():
                return await func(session)

    @property
    def minio_client(self) -> Minio:
        """Lazy-loaded MinIO client."""
        if self._minio_client is None:
            self._minio_client = Minio(
                f"{self.MINIO_HOST}:{self.MINIO_PORT}",
                access_key=self.MINIO_ACCESS_KEY,
                secret_key=self.MINIO_SECRET_KEY,
                secure=False,
            )
        return self._minio_client

    @property
    def qdrant_client(self) -> QdrantClient:
        """Lazy-loaded Qdrant client. Supports self-hosted and Qdrant Cloud."""
        if self._qdrant_client is None:
            use_local = os.getenv("USE_LOCAL_MODE", "false").lower() == "true"
            if use_local:
                self._qdrant_client = QdrantClient(path="./data/qdrant")
            else:
                qdrant_url = os.getenv("QDRANT_URL")
                qdrant_api_key = os.getenv("QDRANT_API_KEY")
                if qdrant_url:
                    self._qdrant_client = QdrantClient(url=qdrant_url, api_key=qdrant_api_key)
                else:
                    self._qdrant_client = QdrantClient(
                        host=self.QDRANT_HOST,
                        port=self.QDRANT_PORT,
                    )
        return self._qdrant_client

    @property
    def embedding_generator(self) -> EmbeddingGenerator:
        """Lazy-loaded EmbeddingGenerator."""
        if self._embedding_generator is None:
            self._embedding_generator = EmbeddingGenerator(self.config_manager)
        return self._embedding_generator

    async def process(self, document_id: uuid.UUID) -> Dict[str, Any]:
        """Fetch, parse, chunk, embed, and index a document by its ID.

        Manages status transitions in PostgreSQL metadata table.

        Parameters
        ----------
        document_id : UUID
            The database UUID document identifier.

        Returns
        -------
        dict
            Contains document_id, status, and metadata about completed execution.
        """
        logger.info("Starting processing pipeline for document: %s", document_id)

        # 1. Update status to 'processing'
        await self._update_db_status(document_id, "processing")

        try:
            # 2. Get metadata from Database
            metadata = await self._fetch_db_metadata(document_id)
            filename = metadata["filename"]
            format_ext = metadata["format"]

            # 3. Read raw bytes from MinIO or local filesystem
            file_bytes = await self._fetch_document(document_id)

            # 4. Parse document text content (enforces 60s timeout)
            parsed_text = await asyncio.wait_for(
                asyncio.to_thread(self._parse_sync, file_bytes, format_ext, filename),
                timeout=60.0
            )

            # 5. Chunk extracted text recursively
            chunks = self._chunk(parsed_text)

            # 6. Generate embeddings and index to Qdrant
            if chunks:
                await self._index_chunks(document_id, chunks)

            # 7. Complete database record status
            await self._update_db_status(document_id, "completed", error_detail=None)
            logger.info("Successfully completed processing for document: %s", document_id)
            return {
                "document_id": document_id,
                "status": "completed",
                "chunks_created": len(chunks)
            }

        except asyncio.TimeoutError as exc:
            err_msg = "Parse failure: timeout after 60 s"
            logger.error("Processing timed out for document %s: %s", document_id, exc)
            await self._update_db_status(document_id, "failed", error_detail=err_msg)
            raise RuntimeError("parse_timeout", err_msg) from exc
        except Exception as exc:
            err_code = "parse_failure"
            err_msg = str(exc)

            if isinstance(exc, EmbeddingFailureError):
                err_code = "embedding_failure"
                err_msg = f"Embedding failure: {exc}"
            elif hasattr(exc, "args") and len(exc.args) > 0 and exc.args[0] in {"ocr_failure", "invalid_file", "parse_failure"}:
                err_code = exc.args[0]
                err_msg = exc.args[1] if len(exc.args) > 1 else str(exc)

            logger.error("Processing failed for document %s: %s", document_id, err_msg)
            await self._update_db_status(document_id, "failed", error_detail=err_msg)
            raise RuntimeError(err_code, err_msg) from exc

    # -----------------------------------------------------------------------
    # Helper Methods
    # -----------------------------------------------------------------------

    async def _fetch_document(self, document_id: uuid.UUID) -> bytes:
        """Retrieve document file content from MinIO object storage or local filesystem."""
        use_local = os.getenv("USE_LOCAL_MODE", "false").lower() == "true"
        if use_local:
            try:
                local_path = os.path.join(os.getcwd(), "data", "documents", str(document_id))
                with open(local_path, "rb") as f:
                    return f.read()
            except Exception as exc:
                logger.error("Failed to read document %s locally: %s", document_id, exc)
                raise RuntimeError("storage_failure", f"Failed to retrieve file locally: {exc}") from exc

        try:
            response = self.minio_client.get_object(self.MINIO_BUCKET, str(document_id))
            file_bytes = response.read()
            response.close()
            response.release_conn()
            return file_bytes
        except Exception as exc:
            logger.error("Failed to read document %s from MinIO: %s", document_id, exc)
            raise RuntimeError("storage_failure", f"Failed to retrieve file from storage: {exc}") from exc

    async def _fetch_db_metadata(self, document_id: uuid.UUID) -> Dict[str, Any]:
        """Fetch filename and format configuration from PostgreSQL database."""
        async def _fetch(session):
            stmt = text(
                """
                SELECT filename, format FROM documents
                WHERE document_id = :document_id
                """
            )
            res = await session.execute(stmt, {"document_id": str(document_id)})
            row = res.fetchone()
            if not row:
                raise RuntimeError("invalid_file", f"Document metadata not found in database for ID: {document_id}")
            return {"filename": row[0], "format": row[1]}

        return await self._execute_db(_fetch)

    async def _update_db_status(self, document_id: uuid.UUID, status: str, error_detail: Optional[str] = None) -> None:
        """Update PostgreSQL processing status for a document."""
        async def _update(session):
            stmt = text(
                """
                UPDATE documents
                SET processing_status = :status, error_detail = :error_detail
                WHERE document_id = :document_id
                """
            )
            await session.execute(stmt, {
                "document_id": str(document_id),
                "status": status,
                "error_detail": error_detail
            })
            await session.commit()

        try:
            await self._execute_db(_update)
        except Exception as exc:
            logger.error("Failed to update status for document %s: %s", document_id, exc)
            raise

    def _parse_sync(self, file_bytes: bytes, format_ext: str, filename: str) -> str:
        """Parse raw file bytes synchronously based on file type extension."""
        logger.debug("Parsing file %s format %s sync", filename, format_ext)
        if not file_bytes:
            raise RuntimeError("invalid_file", "Parse failure: empty file structure")

        try:
            if format_ext == "pdf":
                # Use pypdf — lightweight pure-Python PDF reader
                import pypdf
                reader = pypdf.PdfReader(io.BytesIO(file_bytes))
                pages = []
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        pages.append(text)
                result = "\n\n".join(pages)
                if not result.strip():
                    raise RuntimeError("parse_failure", "Parse failure: PDF contains no extractable text")
                return result

            elif format_ext == "docx":
                # Use python-docx — lightweight DOCX parser
                import docx
                document = docx.Document(io.BytesIO(file_bytes))
                paragraphs = [para.text for para in document.paragraphs if para.text.strip()]
                tables_text = []
                for table in document.tables:
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            tables_text.append(row_text)
                return "\n\n".join(paragraphs + tables_text)

            elif format_ext == "pptx":
                # Use python-pptx — lightweight PPTX parser
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file_bytes))
                slide_texts = []
                for slide in prs.slides:
                    texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            texts.append(shape.text.strip())
                    if texts:
                        slide_texts.append("\n".join(texts))
                return "\n\n".join(slide_texts)

            elif format_ext == "xlsx":
                # Use openpyxl — lightweight XLSX parser
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
                sheets_text = []
                for sheet in wb.worksheets:
                    rows_text = []
                    for row in sheet.iter_rows(values_only=True):
                        row_str = " | ".join(str(cell) for cell in row if cell is not None)
                        if row_str.strip():
                            rows_text.append(row_str)
                    if rows_text:
                        sheets_text.append(f"Sheet: {sheet.title}\n" + "\n".join(rows_text))
                return "\n\n".join(sheets_text)

            elif format_ext == "xls":
                # Use xlrd for legacy XLS files
                import xlrd
                wb = xlrd.open_workbook(file_contents=file_bytes)
                sheets_text = []
                for sheet in wb.sheets():
                    rows_text = []
                    for row_idx in range(sheet.nrows):
                        row_vals = sheet.row_values(row_idx)
                        row_str = " | ".join(str(v) for v in row_vals if str(v).strip())
                        if row_str.strip():
                            rows_text.append(row_str)
                    if rows_text:
                        sheets_text.append(f"Sheet: {sheet.name}\n" + "\n".join(rows_text))
                return "\n\n".join(sheets_text)

            elif format_ext == "txt":
                return file_bytes.decode("utf-8", errors="ignore")

            elif format_ext in {"png", "jpg", "jpeg"}:
                raise RuntimeError("ocr_failure", "OCR engine is not available for image files in this lightweight build.")

            else:
                raise RuntimeError("invalid_file", f"Parse failure: unsupported file format '{format_ext}'")

        except RuntimeError:
            raise
        except Exception as exc:
            logger.error("Sync parsing crashed for %s: %s", filename, exc)
            raise RuntimeError("parse_failure", f"Parse failure: corrupted file structure. Detail: {exc}") from exc

    def _chunk(self, text_content: str) -> List[Dict[str, Any]]:
        """Slice document text content recursively using configured parameters."""
        chunk_size = 400
        overlap_pct = 12.0
        enable_summaries = True

        if self.config_manager:
            try:
                chunk_size = self.config_manager.get("chunking.chunk_size_tokens")
                overlap_pct = self.config_manager.get("chunking.overlap_percentage")
                enable_summaries = self.config_manager.get("chunking.enable_contextual_summaries")
            except Exception as exc:
                logger.warning("Failed to retrieve chunk settings from config: %s", exc)

        overlap_tokens = int(chunk_size * (overlap_pct / 100.0))

        # SentenceSplitter is recursive chunker of LlamaIndex
        splitter = SentenceSplitter(
            chunk_size=chunk_size,
            chunk_overlap=overlap_tokens
        )
        nodes = splitter.get_nodes_from_documents([LlamaDocument(text=text_content)])
        chunks_texts = [node.text for node in nodes]

        processed_chunks = []
        for idx, chunk_text in enumerate(chunks_texts):
            # Compute contextual summary from adjacent text windows
            contextual_summary = ""
            if enable_summaries:
                prev_win = ""
                next_win = ""
                if idx > 0:
                    prev_win = chunks_texts[idx - 1][-150:]
                if idx < len(chunks_texts) - 1:
                    next_win = chunks_texts[idx + 1][:150]
                contextual_summary = f"[Previous context: ...{prev_win}...] [Next context: ...{next_win}...]"

            processed_chunks.append({
                "text": chunk_text,
                "position": idx,
                "contextual_summary": contextual_summary if enable_summaries else None
            })

        return processed_chunks

    async def _index_chunks(self, document_id: uuid.UUID, chunks: List[Dict[str, Any]]) -> None:
        """Create vector embeddings and upsert points to Qdrant vector database."""
        if not chunks:
            return

        # 1. Generate embeddings for all chunks in batches
        points = []
        embedding_size = None
        
        texts = [chunk["text"] for chunk in chunks]
        embeddings = await self.embedding_generator.generate_embeddings(texts)

        for i, chunk in enumerate(chunks):
            embedding = embeddings[i]
            if embedding_size is None:
                embedding_size = len(embedding)

            chunk_id = uuid.uuid4()
            points.append(
                PointStruct(
                    id=str(chunk_id),
                    vector=embedding,
                    payload={
                        "chunk_id": str(chunk_id),
                        "document_id": str(document_id),
                        "chunk_text": chunk["text"],
                        "position_in_document": chunk["position"],
                        "contextual_summary": chunk["contextual_summary"]
                    }
                )
            )

        # 2. Initialize or Recreate Collection in Qdrant to match the embedding size
        try:
            exists = self.qdrant_client.collection_exists(self.QDRANT_COLLECTION)
            if exists:
                info = self.qdrant_client.get_collection(self.QDRANT_COLLECTION)
                existing_size = None
                if hasattr(info, "config") and hasattr(info.config, "params") and hasattr(info.config.params, "vectors"):
                    vectors_config = info.config.params.vectors
                    if hasattr(vectors_config, "size"):
                        existing_size = vectors_config.size
                    elif isinstance(vectors_config, dict) and "size" in vectors_config:
                        existing_size = vectors_config["size"]
                
                if existing_size is not None and existing_size != embedding_size:
                    logger.info("Qdrant collection size mismatch (%d vs %d). Re-creating collection.", existing_size, embedding_size)
                    self.qdrant_client.delete_collection(self.QDRANT_COLLECTION)
                    exists = False

            if not exists:
                self.qdrant_client.create_collection(
                    collection_name=self.QDRANT_COLLECTION,
                    vectors_config=VectorParams(
                        size=embedding_size,
                        distance=Distance.COSINE
                    )
                )
        except Exception as exc:
            logger.error("Failed to check/create Qdrant collection: %s", exc)
            raise RuntimeError("vector_store_failure", f"Qdrant collection check failed: {exc}") from exc

        # 3. Upsert points batch to Qdrant vector store
        try:
            self.qdrant_client.upsert(
                collection_name=self.QDRANT_COLLECTION,
                points=points
            )
        except Exception as exc:
            logger.error("Failed to upsert to Qdrant: %s", exc)
            raise RuntimeError("vector_store_failure", f"Failed to upsert vectors: {exc}") from exc
